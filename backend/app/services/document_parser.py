import os
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation

def extract_text_from_file(file_path):
    """
    Extracts and returns text from a given file path based on its extension.
    Supported extensions: .pdf, .docx, .pptx, .txt
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")
        
    ext = os.path.splitext(file_path)[1].lower()
    
    try:
        if ext == '.pdf':
            return _extract_pdf(file_path)
        elif ext == '.docx':
            return _extract_docx(file_path)
        elif ext == '.pptx':
            return _extract_pptx(file_path)
        elif ext == '.txt':
            return _extract_txt(file_path)
        else:
            raise ValueError(f"Unsupported file format: {ext}")
    except Exception as e:
        raise Exception(f"Failed to extract text from {ext} file: {str(e)}")

def _extract_pdf(file_path):
    text = ""
    with open(file_path, "rb") as f:
        reader = PdfReader(f)
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
    return text.strip()

def _extract_docx(file_path):
    doc = DocxDocument(file_path)
    return "\n".join([para.text for para in doc.paragraphs]).strip()

def _extract_pptx(file_path):
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text.strip()

def _extract_txt(file_path):
    with open(file_path, "r", encoding="utf-8") as f:
        return f.read().strip()
